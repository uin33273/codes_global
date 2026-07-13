"""
Audiobook image-slide generator.

Reproduces the layout rule observed in the reference file
(V1 - コピー.pptx / V1.pptx):

- Normal figure ("図N") slides: picture is placed at a FIXED left
  margin and FIXED width; height scales to preserve the image's
  own aspect ratio.
- QR-code slides: picture keeps a small, roughly fixed box size and
  is centered horizontally.
- The label textbox always shares the same left edge as the picture,
  sits directly above it, and the (label + picture) block is
  vertically centered on the slide by default.

Usage:
    from deck_builder import DeckBuilder
    db = DeckBuilder("template.pptx")   # or DeckBuilder(slide_w=..., slide_h=...)
    db.add_figure_slide("図5", r"image\\img-011.jpg")
    db.add_qr_slide("QRコード2", r"image\\qr2.png")
    db.save("output.pptx")
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from PIL import Image

# --- constants measured from the reference deck (EMU) ---
CONTENT_LEFT = Emu(367025)
CONTENT_WIDTH = Emu(4593600)
LABEL_HEIGHT = Emu(400110)
LABEL_FONT_NAME = "ＭＳ Ｐゴシック"
LABEL_FONT_SIZE = Pt(20)
QR_MAX_SIZE = Emu(1800000)  # ~1.97in box for QR codes


class DeckBuilder:
    def __init__(self, template_path=None, slide_w=None, slide_h=None, layout_index=0):
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        if slide_w:
            self.prs.slide_width = slide_w
        if slide_h:
            self.prs.slide_height = slide_h
        self.layout_index = layout_index

    def _new_slide(self):
        layout = self.prs.slide_layouts[self.layout_index]
        slide = self.prs.slides.add_slide(layout)
        # the reference slides have no leftover "click to add text" placeholders
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)
        return slide

    def _add_label(self, slide, text, left, width=Emu(1400000)):
        box = slide.shapes.add_textbox(left, Emu(0), width, LABEL_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = LABEL_FONT_NAME
        run.font.size = LABEL_FONT_SIZE
        return box

    def add_figure_slide(self, label_text, image_path):
        """Full-width figure: fixed left/width, height keeps aspect ratio."""
        slide = self._new_slide()
        with Image.open(image_path) as im:
            iw, ih = im.size
        pic_width = CONTENT_WIDTH
        pic_height = Emu(int(pic_width * ih / iw))

        block_height = LABEL_HEIGHT + pic_height
        top = Emu(int((self.prs.slide_height - block_height) / 2))

        label = self._add_label(slide, label_text, CONTENT_LEFT)
        label.top = top

        pic = slide.shapes.add_picture(
            image_path, CONTENT_LEFT, Emu(top + LABEL_HEIGHT), pic_width, pic_height
        )
        return slide

    def add_qr_slide(self, label_text, image_path, max_size=QR_MAX_SIZE):
        """Small, horizontally-centered image (QR codes etc.)."""
        slide = self._new_slide()
        with Image.open(image_path) as im:
            iw, ih = im.size
        if iw >= ih:
            pic_width = max_size
            pic_height = Emu(int(max_size * ih / iw))
        else:
            pic_height = max_size
            pic_width = Emu(int(max_size * iw / ih))

        left = Emu(int((self.prs.slide_width - pic_width) / 2))
        block_height = LABEL_HEIGHT + pic_height
        top = Emu(int((self.prs.slide_height - block_height) / 2))

        label = self._add_label(slide, label_text, left, width=pic_width)
        label.top = top

        pic = slide.shapes.add_picture(
            image_path, left, Emu(top + LABEL_HEIGHT), pic_width, pic_height
        )
        return slide

    def save(self, path):
        self.prs.save(path)
