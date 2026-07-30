r"""
Place an image onto a slide, sized by orientation and centered on the
page (both left-right and top-bottom).

Rule:
    landscape image (width >= height): displayed width = LANDSCAPE_WIDTH_CM
    portrait image  (height > width):  displayed height = PORTRAIT_HEIGHT_CM
    (aspect ratio is preserved either way)

Two ways to use:

1) Drag & drop mode (no arguments): double-click this script (or run
   `python place_image_centered.py`), leave the console window open,
   then drag a .pptx file from Explorer onto the console window and
   press Enter. Every slide that already has a picture on it gets
   that picture resized/re-centered per the rule above (the image
   itself is kept, only its size and position are fixed), and the
   "図NN"/"QRコードNN" label is moved to the image's top-left corner.
   The result is saved next to the source file as
   "<元のファイル名>_加工済み.pptx". It then waits for the next file.
   Type nothing and press Enter (or type exit) to quit.

2) One-shot CLI mode (insert a new image into one slide):
   python place_image_centered.py "template.pptx" 1 "image.jpg" "output.pptx"

    template.pptx : the pptx to edit
    1             : 1-based slide number whose picture should be replaced
                    (existing PICTURE shape(s) on that slide are removed)
    image.jpg     : the new image to insert
    output.pptx   : where to save the result (omit to overwrite the input)
"""
import sys
import os
import re
import tempfile
import webbrowser
import tkinter as tk
from tkinter import messagebox, ttk
import numpy as np
from pptx import Presentation
from pptx.util import Cm, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from PIL import Image

LANDSCAPE_WIDTH_CM = 12.76
PORTRAIT_HEIGHT_CM = 16.7
CROP_THRESHOLD = 245  # pixels lighter than this (near-white) are cropped away as margin

# ============================== 仕事の進め方 / 取説 ==============================

INSTRUCTIONS_FILENAME = "手順_PDF画像位置サイズ調整.txt"

DEFAULT_INSTRUCTIONS = """\
※この手順は変更されることがあります。内容を直接書き換えたい場合は
　右上の「✎ 取説を編集」ボタンを押してください(メモ帳でこのファイルが開きます)。
※URL(https://...)やファイルのパス(例: C:\\folder\\file.xlsx)をそのまま1行に
　書いておくと、「仕事の進め方」画面でクリックできるリンクになります。
"""

_LINK_CHAR = r"[^\s。、，」』（）\"']"
LINK_PATTERN = re.compile(rf"(https?://{_LINK_CHAR}+|[A-Za-z]:\\{_LINK_CHAR}+|\\\\{_LINK_CHAR}+)")
LINK_TRAILING_CHARS = "。、,.)）」』\"'"


def get_instructions_path():
    app_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(app_dir, INSTRUCTIONS_FILENAME)


def ensure_instructions_file():
    path = get_instructions_path()
    if not os.path.exists(path):
        with open(path, "w", encoding="utf-8") as f:
            f.write(DEFAULT_INSTRUCTIONS)
    return path


def read_instructions():
    path = ensure_instructions_file()
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except OSError as e:
        return f"手順ファイルの読み込みに失敗しました。\n{path}\n{e}"


def open_instructions_editor():
    path = ensure_instructions_file()
    try:
        os.startfile(path)
    except OSError as e:
        messagebox.showerror("編集に失敗しました", f"{path}\n{e}")


def open_link(href):
    if href.startswith("http://") or href.startswith("https://"):
        webbrowser.open(href)
        return
    try:
        os.startfile(href)
    except OSError as e:
        messagebox.showerror("開けませんでした", f"{href}\n{e}")


def insert_instructions_with_links(text_widget, content):
    text_widget.tag_configure("hyperlink", foreground="#2f6f63", underline=True)
    orig_cursor = text_widget.cget("cursor")
    link_index = 0
    for line in content.splitlines(keepends=True):
        pos = 0
        for m in LINK_PATTERN.finditer(line):
            start, end = m.span()
            if start > pos:
                text_widget.insert("end", line[pos:start])
            raw = m.group(1)
            href = raw.rstrip(LINK_TRAILING_CHARS)
            trailing = raw[len(href):]
            tag_name = f"link_{link_index}"
            link_index += 1
            text_widget.insert("end", href, ("hyperlink", tag_name))
            text_widget.tag_bind(tag_name, "<Button-1>", lambda e, href=href: open_link(href))
            text_widget.tag_bind(tag_name, "<Enter>", lambda e: text_widget.config(cursor="hand2"))
            text_widget.tag_bind(tag_name, "<Leave>", lambda e: text_widget.config(cursor=orig_cursor))
            if trailing:
                text_widget.insert("end", trailing)
            pos = end
        text_widget.insert("end", line[pos:])


def show_instructions_dialog():
    """起動直後に「仕事の進め方」を表示し、閉じるまで先へ進ませない。"""
    root = tk.Tk()
    root.withdraw()

    win = tk.Toplevel(root)
    win.title(f"仕事の進め方({os.path.basename(os.path.dirname(os.path.abspath(__file__)))})")
    win.geometry("560x520")

    text_frame = ttk.Frame(win)
    text_frame.pack(fill="both", expand=True, padx=14, pady=(14, 6))

    text_widget = tk.Text(text_frame, wrap="word")
    scroll = ttk.Scrollbar(text_frame, orient="vertical", command=text_widget.yview)
    text_widget.config(yscrollcommand=scroll.set)
    scroll.pack(side="right", fill="y")
    text_widget.pack(side="left", fill="both", expand=True)

    insert_instructions_with_links(text_widget, read_instructions())
    text_widget.config(state="disabled")

    button_row = ttk.Frame(win)
    button_row.pack(side="bottom", fill="x", padx=14, pady=(0, 14))
    ttk.Button(button_row, text="✎ 取説を編集", command=open_instructions_editor).pack(side="left")
    ttk.Button(button_row, text="閉じる(作業を始める)", command=win.destroy).pack(side="right")

    win.protocol("WM_DELETE_WINDOW", win.destroy)
    win.attributes("-topmost", True)
    win.lift()
    win.focus_force()
    root.wait_window(win)
    root.destroy()

# textboxes like "図00", "図1", "図2-1", "QRコード1" are treated as the image's label
LABEL_PATTERN = re.compile(r"^(図|QRコード)\s*[0-9０-９]+(-[0-9０-９]+)*$")


def is_picture_shape(shape):
    """True for any <p:pic> element, including pictures that were inserted
    into a placeholder (python-pptx reports those as shape_type PLACEHOLDER,
    not PICTURE, even though they hold real image data)."""
    return shape._element.tag == qn("p:pic")


def find_label_shape(slide):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
            if LABEL_PATTERN.match(shape.text_frame.text.strip()):
                return shape
    return None


def autocrop(im, threshold=CROP_THRESHOLD):
    """Crop away near-white whitespace margins, keeping the image's own mode."""
    rgb = im.convert("RGB")
    arr = np.array(rgb)
    # a pixel counts as "content" if any channel is darker than threshold
    mask = np.any(arr < threshold, axis=2)
    rows = np.where(mask.any(axis=1))[0]
    cols = np.where(mask.any(axis=0))[0]
    if rows.size == 0 or cols.size == 0:
        return im  # nothing but background found; return unchanged
    top, bottom = rows[0], rows[-1]
    left, right = cols[0], cols[-1]
    return im.crop((left, top, right + 1, bottom + 1))


def place_image_centered(slide, slide_width, slide_height, image_path):
    # remove existing pictures on the slide so the new one takes their place
    for shape in list(slide.shapes):
        if is_picture_shape(shape):
            shape._element.getparent().remove(shape._element)

    # crop whitespace margins first, then size/center based on the
    # cropped image so the label is positioned against the final bounds
    ext = os.path.splitext(image_path)[1] or ".png"
    with Image.open(image_path) as im:
        cropped = autocrop(im)
        iw, ih = cropped.size
        fd, cropped_path = tempfile.mkstemp(suffix=ext)
        os.close(fd)
        cropped.save(cropped_path)

    try:
        if iw >= ih:  # landscape (or square)
            width = Cm(LANDSCAPE_WIDTH_CM)
            height = Emu(int(width * ih / iw))
        else:  # portrait
            height = Cm(PORTRAIT_HEIGHT_CM)
            width = Emu(int(height * iw / ih))
            # a near-square "portrait" image can end up wider than the page;
            # cap the width and derive height from that instead
            if width > Cm(LANDSCAPE_WIDTH_CM):
                width = Cm(LANDSCAPE_WIDTH_CM)
                height = Emu(int(width * ih / iw))

        left = Emu(int((slide_width - width) / 2))
        top = Emu(int((slide_height - height) / 2))

        slide.shapes.add_picture(cropped_path, left, top, width, height)
    finally:
        os.remove(cropped_path)

    # align the "図NN" label's left edge with the image's left edge, and
    # its bottom edge with the image's top edge (label directly above the
    # image, text stays visible), wherever it happened to be placed originally
    label = find_label_shape(slide)
    if label is not None:
        label.left = left
        label.top = Emu(int(top - label.height))


def fix_existing_pictures(prs):
    """Re-size/re-center every picture already present in the deck,
    keeping the image itself. Returns the number of pictures fixed."""
    n = 0
    with tempfile.TemporaryDirectory() as tmp_dir:
        for slide in prs.slides:
            pics = [s for s in slide.shapes if is_picture_shape(s)]
            for pic in pics:
                tmp_path = os.path.join(tmp_dir, f"_extract_{n}.{pic.image.ext}")
                with open(tmp_path, "wb") as f:
                    f.write(pic.image.blob)
                place_image_centered(slide, prs.slide_width, prs.slide_height, tmp_path)
                n += 1
    return n


def split_paths(line):
    """Split a line of one or more (optionally quoted) Windows paths
    without mangling backslashes the way shlex.split() would."""
    line = line.strip()
    # PowerShell prefixes a dragged file/folder with the call operator
    # "&" and single-quotes the path (e.g. & 'C:\...\file.pptx'); strip
    # that operator so the quoted path below is parsed normally
    if line.startswith("&"):
        line = line[1:].strip()
    return [
        dq or sq or bare
        for dq, sq, bare in re.findall(r'"([^"]+)"|\'([^\']+)\'|(\S+)', line)
    ]


def process_one(path):
    path = path.strip().strip('"').strip("'")
    if not path:
        return None
    if not os.path.isfile(path):
        print(f"  !! ファイルが見つかりません: {path}")
        return None
    if not path.lower().endswith(".pptx"):
        print(f"  !! .pptx ではありません: {path}")
        return None

    try:
        prs = Presentation(path)
        n = fix_existing_pictures(prs)

        base = os.path.splitext(os.path.basename(path))[0]
        output_dir = os.path.dirname(os.path.abspath(path))
        output_path = os.path.join(output_dir, f"{base}_加工済み.pptx")
        prs.save(output_path)

        print(f"  {n}枚の画像を調整しました")
        print(f"  保存先: {output_path}")
        return output_path
    except Exception as e:
        print(f"  !! エラー: {e}")
        return None


def interactive_loop():
    print("=== PDF画像調整 (ドラッグ&ドロップ待機モード) ===")
    print("保存先フォルダ: ドラッグしたファイルと同じフォルダ")
    print(".pptx ファイルをこのウィンドウにドラッグ&ドロップして Enter を押してください。")
    print("終了するには何も入力せず Enter、または exit と入力してください。")
    print()
    while True:
        try:
            line = input("> ").strip()
        except EOFError:
            break
        if not line or line.lower() in ("exit", "quit"):
            print("終了します。")
            break
        output_paths = [p for p in (process_one(path) for path in split_paths(line)) if p]
        print("完了しました。")
        for output_path in output_paths:
            os.startfile(output_path)
        return


def main():
    show_instructions_dialog()
    if len(sys.argv) == 1:
        interactive_loop()
        return

    # a single .pptx argument (e.g. dragged onto the .py/.bat icon directly)
    # is handled the same way as a line typed into the interactive loop
    if len(sys.argv) == 2 and sys.argv[1].lower().endswith(".pptx"):
        output_path = process_one(sys.argv[1])
        if output_path:
            os.startfile(output_path)
        return

    template_path = sys.argv[1]
    slide_no = int(sys.argv[2])
    image_path = sys.argv[3]
    output_path = sys.argv[4] if len(sys.argv) > 4 else template_path

    prs = Presentation(template_path)
    slide = list(prs.slides)[slide_no - 1]

    place_image_centered(slide, prs.slide_width, prs.slide_height, image_path)

    prs.save(output_path)
    print("saved:", output_path)


if __name__ == "__main__":
    main()
